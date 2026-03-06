#!/usr/bin/env python3
"""Simple generator: create proposal doc, PPT summary, pricing CSV, and Gantt CSV from parsed RFP JSON."""
import sys
import json
import csv
from pathlib import Path
import re

# When this script is executed as `python proposal_agent/tools/generator.py`
# the package `proposal_agent` may not be on `sys.path`. Insert the repo root
# so `from proposal_agent.kb import ...` works both under `python -m pytest`
# and when the script is run directly as a subprocess in tests.
try:
    repo_root = Path(__file__).resolve().parents[2]
    sys.path.insert(0, str(repo_root))
except Exception:
    pass

# Use the KB loader for clause ingestion/search
from proposal_agent.kb import loader as kb_loader


def _parse_markdown_to_blocks(text: str):
    """Very small markdown parser: returns list of blocks (type, content).

    Types: 'heading', 'para', 'bullets' (list of lines)
    """
    lines = [l.rstrip() for l in text.splitlines()]
    blocks = []
    buf = []
    bullets = []

    def flush_para():
        nonlocal buf
        if buf:
            blocks.append(('para', '\n'.join(buf).strip()))
            buf = []

    def flush_bullets():
        nonlocal bullets
        if bullets:
            blocks.append(('bullets', bullets))
            bullets = []

    for line in lines:
        if not line.strip():
            flush_para()
            flush_bullets()
            continue
        if line.lstrip().startswith('#'):
            flush_para()
            flush_bullets()
            # heading level ignored, just use text after hashes
            h = line.lstrip().lstrip('#').strip()
            blocks.append(('heading', h))
            continue
        if line.lstrip().startswith('- '):
            bullets.append(line.lstrip()[2:].strip())
            continue
        # normal paragraph line
        buf.append(line)

    flush_para()
    flush_bullets()
    return blocks


def _render_blocks_into_doc(doc, blocks):
    for typ, content in blocks:
        if typ == 'heading':
            doc.add_heading(content, level=3)
        elif typ == 'para':
            for para in content.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())
        elif typ == 'bullets':
            for b in content:
                doc.add_paragraph(b, style='List Bullet')


def _render_blocks_into_ppt(text_frame, blocks):
    # Ensure at least one paragraph exists
    first = True
    for typ, content in blocks:
        if typ == 'heading':
            if first:
                text_frame.text = content
                first = False
            else:
                p = text_frame.add_paragraph()
                p.text = content
                p.level = 0
        elif typ == 'para':
            for para in content.split('\n\n'):
                if para.strip():
                    p = text_frame.add_paragraph()
                    p.text = para.strip()
                    p.level = 1
        elif typ == 'bullets':
            for b in content:
                p = text_frame.add_paragraph()
                p.text = b
                p.level = 1


# Optional rich exports
try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

try:
    import openpyxl
    from openpyxl import Workbook
except Exception:
    openpyxl = None



def load_rates(path: Path):
    rates = {}
    with path.open() as f:
        reader = csv.DictReader(f)
        for r in reader:
            role = r['role']
            rates[role] = {
                'onshore': float(r['onshore_daily_usd']),
                'offshore': float(r['offshore_daily_usd']),
                'util': float(r['default_utilization_pct'])/100.0
            }
    return rates


def estimate_price(parsed: dict, rates: dict, onshore_frac: float = 0.4, contingency_frac: float = 0.2, days_overrides: dict = None):
    # Load structured estimation model if available
    model_path = Path(__file__).parents[1] / 'data' / 'estimation_model.json'
    model = {}
    if model_path.exists():
        try:
            model = json.loads(model_path.read_text(encoding='utf-8'))
        except Exception:
            model = {}

    modules = max(1, int(parsed.get('modules_count', 1)))

    # Determine product profile from title or parsed product
    title_l = parsed.get('title', '').lower()
    profile_key = None
    if 's/4' in title_l or 'sap' in title_l:
        profile_key = 'sap'
    elif 'dynamics' in title_l or 'd365' in title_l or 'dynamics 365' in title_l:
        profile_key = 'dynamics'
    elif 'mes' in title_l:
        profile_key = 'mes'
    # allow explicit product in parsed
    if not profile_key and parsed.get('product'):
        profile_key = parsed.get('product').lower()

    days_per_module = model.get('default_days_per_module', 10)
    module_multiplier = 1.0
    mix = None
    if model and profile_key and model.get('product_profiles', {}).get(profile_key):
        prof = model['product_profiles'][profile_key]
        days_per_module = prof.get('days_per_module', days_per_module)
        module_multiplier = prof.get('module_multiplier', 1.0)
        mix = prof.get('mix')

    # apply product-level calibration if available, otherwise global
    product_cal_map = model.get('product_calibration', {}) if model else {}
    global_cal = float(model.get('calibration_factor', 1.0)) if model else 1.0
    prod_cal = product_cal_map.get(profile_key, global_cal)
    total_days = modules * days_per_module * module_multiplier * float(prod_cal)

    # apply customer adjustment multiplier if available
    cust_adj_path = Path(__file__).parents[1] / 'data' / 'customer_adjustments.json'
    cust_mult = 1.0
    if cust_adj_path.exists():
        try:
            cadj = json.loads(cust_adj_path.read_text(encoding='utf-8'))
            company = parsed.get('company','')
            if company and company in cadj:
                cust_mult = float(cadj[company])
            else:
                cust_mult = float(cadj.get('default', 1.0))
        except Exception:
            cust_mult = 1.0
    total_days = total_days * cust_mult

    # Fallback mix if model not present
    if not mix:
        mix = {
            'Project Manager': 0.10,
            'Solution Architect': 0.10,
            'Functional Consultant': 0.30,
            'Developer': 0.30,
            'QA Engineer': 0.15,
            'Basis/Infra': 0.05
        }

    staffing = []
    total_price = 0.0

    for role, pct in mix.items():
        # calculate base days then apply overrides
        base_days = max(1, int(round(total_days * pct)))
        if days_overrides and role in days_overrides:
            try:
                role_days = int(days_overrides[role])
            except Exception:
                role_days = base_days
        else:
            role_days = base_days

        # get rates: prefer rates param, otherwise model defaults
        r = None
        if rates and role in rates:
            r = rates.get(role)
        elif model.get('roles') and role in model['roles']:
            rr = model['roles'][role]
            r = {'onshore': rr.get('onshore'), 'offshore': rr.get('offshore'), 'util': rr.get('util', 0.8)}

        if not r:
            daily = 700.0
        else:
            daily = onshore_frac * float(r['onshore']) + (1.0 - onshore_frac) * float(r['offshore'])

        price = role_days * daily
        staffing.append({'role': role, 'days': role_days, 'daily_rate': round(daily, 2), 'price': round(price, 2)})
        total_price += price

    contingency = contingency_frac * total_price
    fixed_price = total_price + contingency
    return staffing, round(total_price, 2), round(contingency, 2), round(fixed_price, 2)


def write_proposal_doc(parsed: dict, staffing, fixed_price: float, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    # Prefer creating a real .docx when python-docx is installed
    doc_path = outdir / 'proposal.docx'
    if Document:
        doc = Document()
        doc.add_heading(parsed.get('title', 'Proposal'), level=1)
        doc.add_heading('Executive Summary', level=2)
        doc.add_paragraph('This proposal provides a fixed-price approach for the requested scope.')

        doc.add_heading('Scope', level=2)
        for s in parsed.get('scope_items', []):
            doc.add_paragraph(s, style='List Bullet')

        doc.add_heading('Staffing & Pricing', level=2)
        table = doc.add_table(rows=1, cols=4)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Role'
        hdr_cells[1].text = 'Days'
        hdr_cells[2].text = 'Daily Rate (USD)'
        hdr_cells[3].text = 'Price (USD)'
        for s in staffing:
            row_cells = table.add_row().cells
            row_cells[0].text = s['role']
            row_cells[1].text = str(s['days'])
            row_cells[2].text = f"{s['daily_rate']:.2f}"
            row_cells[3].text = f"{s['price']:.2f}"

        doc.add_paragraph(f"Subtotal: ${sum(s['price'] for s in staffing):.2f}")
        doc.add_paragraph(f"Contingency (20%): ${0.2 * sum(s['price'] for s in staffing):.2f}")
        doc.add_paragraph(f"Fixed Price (Total): ${fixed_price:.2f}")

        # Insert detailed SOW from templates (insert before Staffing & Pricing)
        template_dir = Path(__file__).parents[1] / 'kb' / 'templates'
        templates = load_clauses(template_dir) if False else None
        # load templates differently (read md files)
        templates = {}
        for p in (template_dir).glob('*.md'):
            templates[p.stem.lower()] = p.read_text(encoding='utf-8')

        # Determine product type
        title_l = parsed.get('title','').lower()
        chosen = None
        if 's/4' in title_l or 'sap' in title_l:
            chosen = templates.get('sap_s4_template') or templates.get('sap_s4')
        elif 'dynamics' in title_l or 'd365' in title_l or 'dynamics 365' in title_l:
            chosen = templates.get('dynamics_template') or templates.get('dynamics')
        elif 'mes' in title_l:
            chosen = templates.get('mes_template') or templates.get('mes')

        def render_template(text: str, parsed: dict) -> str:
            # Replace simple placeholders: {{company}}, {{submission_deadline}}, {{scope_items}}
            out = text
            out = out.replace('{{company}}', parsed.get('company', ''))
            out = out.replace('{{submission_deadline}}', parsed.get('submission_deadline', ''))
            # scope_items placeholder: render as bullet list
            if '{{scope_items}}' in out:
                items = parsed.get('scope_items', [])
                if items:
                    bullets = '\n'.join([f"- {i}" for i in items])
                else:
                    bullets = '- [Scope details here]'
                out = out.replace('{{scope_items}}', bullets)
            return out
        def distribute_scope_items(template_text: str, scope_items: list) -> str:
            # Greedy keyword mapping of scope items into template subsections.
            # Returns modified template text with injected bullets under matching headings.
            out = template_text
            # Define mapping keywords -> heading labels in template
            mappings = {
                'finance': ['finance', 'ledger', 'accounts', 'ar/', 'ap', 'asset'],
                'mm': ['materials', 'mm', 'inventory', 'warehouse'],
                'sd': ['sales', 'order', 'billing', 'pricing', 'sd'],
                'pp': ['production', 'pp', 'planning', 'production orders'],
                'bw_bi': ['bw', 'bi', 'analytics', 'report', 'dashboard'],
                'mes': ['mes', 'shop-floor', 'plc', 'oee', 'traceability'],
                'integration': ['integrat', 'api', 'plc', 'erp', 'connect'],
                'migration': ['migrat', 'data', 'historical', 'convert'],
            }

            # Locate headings by scanning lines; build insertion map
            insert_map = {}
            lines = out.splitlines()
            for i, line in enumerate(lines):
                # find headings like 'Finance:' or lines starting with '- Finance:'
                m = re.match(r"^(#+\s*)?(?P<h>[^:\n]+):\s*$", line.strip())
                if m:
                    heading = m.group('h').strip().lower()
                    insert_map[heading] = i

            # For each scope item, assign to a heading if keyword matches
            extras = {k: [] for k in insert_map.keys()}
            unassigned = []
            for item in scope_items:
                assigned = False
                low = item.lower()
                for h, kws in mappings.items():
                    for kw in kws:
                        if kw in low:
                            # try to map to a heading that best matches h
                            # prefer exact headings like 'finance' or 'mm'
                            for heading in insert_map.keys():
                                if h in heading:
                                    extras[heading].append(item)
                                    assigned = True
                                    break
                            if assigned:
                                break
                    if assigned:
                        break
                if not assigned:
                    unassigned.append(item)

            # Inject bullets under headings by reconstructing lines
            new_lines = []
            for i, line in enumerate(lines):
                new_lines.append(line)
                # check if this line is a heading we should inject under
                key = line.strip().lower().rstrip(':')
                if key in extras and extras[key]:
                    for it in extras[key]:
                        new_lines.append(f"- {it}")

            # Append unassigned items near top under a 'Customer Scope Items' marker if present
            joined = '\n'.join(new_lines)
            if '{{scope_items}}' in joined:
                bullets = '\n'.join([f"- {u}" for u in unassigned]) if unassigned else ''
                joined = joined.replace('{{scope_items}}', bullets or '- [Scope details here]')
            else:
                if unassigned:
                    joined += '\n\nCustomer Scope Items:\n' + '\n'.join([f"- {u}" for u in unassigned])
            return joined

        if chosen:
            doc.add_page_break()
            doc.add_heading('Detailed Statement of Work (SOW)', level=2)
            # render placeholders first
            rendered = render_template(chosen, parsed)
            # then try to distribute scope items into subsection bullets
            rendered2 = distribute_scope_items(rendered, parsed.get('scope_items', []))
            for para in rendered2.split('\n\n'):
                if para.strip():
                    if para.strip().startswith('#'):
                        doc.add_heading(para.strip('# ').strip(), level=3)
                    else:
                        # write each line separately to preserve bullets
                        for line in para.splitlines():
                            if line.strip().startswith('-'):
                                doc.add_paragraph(line.strip('- ').strip(), style='List Bullet')
                            else:
                                doc.add_paragraph(line.strip())

        # Insert compliance/legal clauses if present in KB
        clause_dir = Path(__file__).parents[1] / 'kb' / 'clauses'
        clauses = load_clauses(clause_dir)
        selected = select_clauses(parsed, clauses)
        if selected:
            doc.add_page_break()
            doc.add_heading('Compliance & Legal Clauses', level=2)
            for title, body in selected:
                doc.add_heading(title, level=3)
                blocks = _parse_markdown_to_blocks(body)
                # If the clause has no explicit blocks, treat whole as paragraph
                if not blocks:
                    doc.add_paragraph(body)
                else:
                    _render_blocks_into_doc(doc, blocks)

        doc.save(doc_path)
        return doc_path
    else:
        # Fallback to plain text file named .docx
        with doc_path.open('w', encoding='utf-8') as f:
            f.write(parsed.get('title','Proposal') + '\n\n')
            f.write('Executive Summary:\n')
            f.write('This proposal provides a fixed-price approach for the requested scope.\n\n')
            f.write('Scope:\n')
            for s in parsed.get('scope_items',[]):
                f.write(f"- {s}\n")
            f.write('\nStaffing & Pricing:\n')
            for s in staffing:
                f.write(f"{s['role']}: {s['days']} days @ ${s['daily_rate']:.2f}/day = ${s['price']:.2f}\n")
            f.write(f"\nSubtotal: ${sum(s['price'] for s in staffing):.2f}\n")
            f.write(f"Contingency (20%): ${0.2 * sum(s['price'] for s in staffing):.2f}\n")
            f.write(f"Fixed Price (Total): ${fixed_price:.2f}\n")
            # Append compliance clauses to fallback text
            clause_dir = Path(__file__).parents[1] / 'kb' / 'clauses'
            clauses = load_clauses(clause_dir)
            selected = select_clauses(parsed, clauses)
            if selected:
                f.write('\nCompliance & Legal Clauses:\n')
                for title, body in selected:
                    f.write('\n' + title + '\n')
                    # simplify markdown to plain text
                    lines = [l for l in body.splitlines() if not l.strip().startswith('```')]
                    for line in lines:
                        f.write(line + '\n')
        return doc_path


def write_proposal_ppt(parsed: dict, staffing, fixed_price: float, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    ppt_path = outdir / 'exec_summary.pptx'
    if Presentation:
        prs = Presentation()
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = parsed.get('title', 'Proposal')
        subtitle = slide.placeholders[1]
        subtitle.text = 'Executive Summary'

        # Scope slide
        body_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(body_slide_layout)
        slide.shapes.title.text = 'Scope'
        tf = slide.shapes.placeholders[1].text_frame
        for s in parsed.get('scope_items', []):
            p = tf.add_paragraph()
            p.text = s
            p.level = 1

        # Pricing slide
        slide = prs.slides.add_slide(body_slide_layout)
        slide.shapes.title.text = 'Pricing Snapshot'
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = f"Fixed Price: ${fixed_price:.2f}"

        # add compliance slide if clauses
        clause_dir = Path(__file__).parents[1] / 'kb' / 'clauses'
        clauses = load_clauses(clause_dir)
        selected = select_clauses(parsed, clauses)
        if selected:
            slide = prs.slides.add_slide(body_slide_layout)
            slide.shapes.title.text = 'Compliance & Legal Clauses'
            tf = slide.shapes.placeholders[1].text_frame
            for title, body in selected:
                p_title = tf.add_paragraph()
                p_title.text = title
                p_title.level = 0
                blocks = _parse_markdown_to_blocks(body)
                if blocks:
                    _render_blocks_into_ppt(tf, blocks)
                else:
                    first_para = body.split('\n\n')[0].strip()
                    p = tf.add_paragraph()
                    p.text = first_para[:200]
                    p.level = 1

        prs.save(ppt_path)
        return ppt_path
    else:
        # Fallback plain text summary saved with .pptx extension
        with ppt_path.open('w', encoding='utf-8') as f:
            f.write(parsed.get('title','Proposal') + '\n\n')
            f.write('Executive Summary (PPT fallback)\n')
            for s in parsed.get('scope_items',[]):
                f.write(f"- {s}\n")
            f.write(f"\nFixed Price: ${fixed_price:.2f}\n")
        return ppt_path


def write_pricing_csv(staffing, outdir: Path):
    outdir.mkdir(parents=True, exist_ok=True)
    path = outdir / 'pricing.csv'
    with path.open('w', newline='', encoding='utf-8') as f:
        w = csv.writer(f)
        w.writerow(['role','days','daily_rate','price'])
        for s in staffing:
            w.writerow([s['role'], s['days'], f"{s['daily_rate']:.2f}", f"{s['price']:.2f}"])
    return path


def write_pricing_xlsx(staffing, outdir: Path):
    path = outdir / 'pricing.xlsx'
    if openpyxl:
        wb = Workbook()
        ws = wb.active
        ws.title = 'Pricing'
        ws.append(['Role','Days','Daily Rate (USD)','Price (USD)'])
        for s in staffing:
            ws.append([s['role'], s['days'], round(s['daily_rate'],2), round(s['price'],2)])
        # Totals and contingency formulas
        last_row = ws.max_row + 1
        ws.append(['Subtotal', '', '', f"=SUM(D2:D{last_row-1})"]) 
        ws.append(['Contingency (20%)', '', '', f"=D{last_row}*0.2"]) 
        ws.append(['Fixed Price (Total)', '', '', f"=D{last_row}+D{last_row+1}"])
        wb.save(path)
        return path
    else:
        # Fallback to CSV
        return write_pricing_csv(staffing, outdir)


def write_gantt(parsed: dict, outdir: Path):
    path = outdir / 'gantt.csv'
    # Very simple gantt: milestone per module, 1 week each
    modules = parsed.get('scope_items', [])
    with path.open('w', newline='', encoding='utf-8') as f:
        w = csv.writer(f)
        w.writerow(['task','start_week','duration_weeks'])
        for i, m in enumerate(modules, start=1):
            w.writerow([m[:60], i, 1])
    return path


def load_clauses(clause_dir: Path) -> dict:
    # Delegate to proposal_agent.kb.loader which understands the KB layout
    # Accept either the `kb` root or the `kb/clauses` folder.
    if clause_dir.name == 'clauses':
        kb_root = clause_dir.parent
    else:
        kb_root = clause_dir
    try:
        return kb_loader.load_clauses(kb_root)
    except Exception:
        # Fall back to simple reader
        clauses = {}
        if not clause_dir.exists():
            return clauses
        for p in clause_dir.glob('*.md'):
            key = p.stem.lower()
            clauses[key] = p.read_text(encoding='utf-8')
        return clauses


def select_clauses(parsed: dict, clauses: dict) -> list:
    selected = []
    comp = parsed.get('compliance', '') or ''
    comp_l = comp.lower()

    # Direct mapping of known clause keys
    mapping = [
        ('gdpr', 'GDPR', 'GDPR'),
        ('iso', 'ISO', 'ISO'),
        ('data_residency', 'Data Residency', 'Data Residency'),
        ('local_procurement', 'Local Procurement', 'Local Procurement'),
    ]
    def _text_for(key):
        v = clauses.get(key)
        if v is None:
            return ''
        if isinstance(v, dict):
            return v.get('text','')
        return v

    for key, title, _ in mapping:
        if key in clauses and (key in comp_l or key.replace('_', ' ') in comp_l):
            selected.append((title, _text_for(key)))

    # If compliance string contains keywords not directly mapped, try a keyword search
    keywords = []
    if 'gdpr' in comp_l or 'data protection' in comp_l:
        keywords.append('gdpr')
    if 'iso' in comp_l or 'security' in comp_l:
        keywords.append('iso')
    if 'data residency' in comp_l or 'residency' in comp_l:
        keywords.append('data residency')
    if 'procure' in comp_l or 'local procurement' in comp_l:
        keywords.append('procurement')

    for kw in keywords:
        matches = kb_loader.find_clauses_by_keyword(clauses, kw)
        for m in matches:
            # avoid duplicates
            if any(m == s[0].lower().replace(' ', '_') for s in selected):
                continue
            title = m.replace('_', ' ').title()
            selected.append((title, _text_for(m)))

    # Tag-based selection: if parsed contains `clause_tags` (list or comma string), include clauses with matching meta.tags
    tag_sources = parsed.get('clause_tags') or parsed.get('required_clause_tags') or parsed.get('clause_tags') or parsed.get('tags')
    if tag_sources:
        if isinstance(tag_sources, str):
            tag_list = [t.strip() for t in tag_sources.split(',') if t.strip()]
        else:
            tag_list = list(tag_sources)

        for name, val in clauses.items():
            meta_tags = []
            if isinstance(val, dict):
                meta_tags = [t.lower() for t in val.get('meta', {}).get('tags', [])]
            # fallback: if clause text contains tag words, include
            if any(t.lower() in meta_tags for t in tag_list):
                if any(name == s[0].lower().replace(' ', '_') for s in selected):
                    continue
                selected.append((name.replace('_', ' ').title(), _text_for(name)))

    return selected


if __name__ == '__main__':
    import argparse

    p = argparse.ArgumentParser(description='Generate proposal artifacts from parsed RFP JSON')
    p.add_argument('parsed_json', help='Parsed RFP JSON file')
    p.add_argument('--rates-file', help='CSV file with role rates (overrides default)', default=None)
    # Escape percent signs in help strings so argparse doesn't attempt
    # to interpret them as formatting placeholders on some Python versions.
    p.add_argument('--contingency', type=float, help='Contingency as fraction (e.g. 0.2 for 20%%)', default=0.20)
    p.add_argument('--onshore-pct', type=float, help='Onshore percentage for blended rates (0-100)', default=40.0)
    p.add_argument('--outdir', help='Output directory (defaults to parsed JSON parent/output)', default=None)
    p.add_argument('--set-rate', action='append', help='Override role rate. Formats: "Role=500" or "Role:onshore=600" or "Role:offshore=300". Can be repeated.', default=[])
    p.add_argument('--set-days', action='append', help='Override role effort days. Format: "Role=30". Can be repeated.', default=[])
    args = p.parse_args()

    parsed_path = Path(args.parsed_json)
    if not parsed_path.exists():
        print('Parsed JSON not found:', parsed_path)
        sys.exit(1)

    parsed = json.loads(parsed_path.read_text(encoding='utf-8'))

    rates_path = Path(args.rates_file) if args.rates_file else (Path(__file__).parents[1] / 'data' / 'default_rates.csv')
    if not rates_path.exists():
        print('Rates file not found:', rates_path)
        rates = {}
    else:
        rates = load_rates(rates_path)

    # Apply per-role overrides from CLI
    for override in args.set_rate:
        # support Role=rate or Role:onshore=rate or Role:offshore=rate
        try:
            if ':' in override and '=' in override:
                # Role:onshore=600
                left, val = override.split('=', 1)
                role, which = left.split(':', 1)
                role = role.strip()
                which = which.strip().lower()
                rate_val = float(val.strip())
                if role not in rates:
                    rates[role] = {'onshore': rate_val, 'offshore': rate_val, 'util': 0.8}
                else:
                    if which == 'onshore':
                        rates[role]['onshore'] = rate_val
                    elif which == 'offshore':
                        rates[role]['offshore'] = rate_val
            elif '=' in override:
                # Role=500 -> set both onshore and offshore to 500
                role, val = override.split('=', 1)
                role = role.strip()
                rate_val = float(val.strip())
                rates[role] = {'onshore': rate_val, 'offshore': rate_val, 'util': 0.8}
        except Exception as e:
            print(f"Warning: couldn't parse set-rate '{override}': {e}")

    # Convert percent to fraction
    onshore_frac = max(0.0, min(100.0, args.onshore_pct)) / 100.0

    # parse --set-days overrides
    days_overrides = {}
    if hasattr(args, 'set_days') and args.set_days:
        for ov in args.set_days:
            try:
                if '=' in ov:
                    role, val = ov.split('=', 1)
                    days_overrides[role.strip()] = int(val.strip())
            except Exception as e:
                print(f"Warning: couldn't parse set-days '{ov}': {e}")

    staffing, subtotal, contingency_amount, fixed_price = estimate_price(parsed, rates, onshore_frac, args.contingency, days_overrides)

    outdir = Path(args.outdir) if args.outdir else (parsed_path.parent / 'output')
    doc = write_proposal_doc(parsed, staffing, fixed_price, outdir)
    ppt = write_proposal_ppt(parsed, staffing, fixed_price, outdir)
    pricing_csv = write_pricing_csv(staffing, outdir)
    pricing_xlsx = write_pricing_xlsx(staffing, outdir)
    gantt = write_gantt(parsed, outdir)
    print('Generated:', doc, ppt, pricing_csv, pricing_xlsx, gantt)
